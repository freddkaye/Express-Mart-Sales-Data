from pptx import Presentation

# Create a new PowerPoint presentation
presentation = Presentation()

# Add a title slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "ExpressMart Sales Analysis Report"
subtitle.text = "Enhancing Business Performance through Data Insights"

# Add an Objective slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objective"
content.text = (
    "The goal of this analysis is to enhance business performance by:\n"
    "- Understanding customer purchasing behaviors\n"
    "- Identifying sales trends\n"
    "- Optimizing operations across locations"
)

# Add a slide for Peak Sales Periods
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Peak Sales Periods"
content.text = (
    "Findings:\n"
    "- Peak hours: Mornings and evenings.\n"
    "- Peak months: Driven by seasonal demand.\n\n"
    "Recommendations:\n"
    "- Optimize staffing during peak times.\n"
    "- Increase inventory for high-demand products."
)

# Add a slide for Top Performing Products
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Top Performing Products"
content.text = (
    "Findings:\n"
    "- Top sellers: Product A and Product B.\n"
    "- Underperforming: Product X contributes minimally to revenue.\n\n"
    "Recommendations:\n"
    "- Promote underperforming products with targeted campaigns.\n"
    "- Bundle top sellers with other items for cross-selling."
)

# Add a slide for Regional Sales Performance
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Regional Sales Performance"
content.text = (
    "Findings:\n"
    "- Leading cities: New York City and San Francisco.\n"
    "- Growth potential: Atlanta and Austin.\n\n"
    "Recommendations:\n"
    "- Launch localized marketing campaigns in underperforming regions.\n"
    "- Offer region-specific discounts and offers."
)

# Add a slide for Customer Purchasing Behaviors
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Customer Purchasing Behaviors"
content.text = (
    "Findings:\n"
    "- Frequent purchases during weekends.\n"
    "- Commonly paired items: Item 1 and Item 2.\n\n"
    "Recommendations:\n"
    "- Promote weekend-exclusive deals.\n"
    "- Tailor marketing based on demographics and purchasing patterns."
)

# Add a slide for Sales Trends Analysis
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Sales Trends Analysis"
content.text = (
    "Findings:\n"
    "- Seasonal variations: Q4 outperforms due to holidays.\n"
    "- Low-performing periods: Summer months.\n\n"
    "Recommendations:\n"
    "- Run summer promotions to boost sales.\n"
    "- Expand product lines for seasonal demand shifts."
)

# Add a slide for the Sales Performance Dashboard
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Sales Performance Dashboard"
content.text = (
    "Deliverable:\n"
    "- Dynamic Excel dashboard for tracking KPIs (total sales, quantity sold, etc.).\n\n"
    "Purpose:\n"
    "- Provides real-time insights for data-driven decisions."
)

# Add a Conclusion slide
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "This analysis provides actionable insights to:\n"
    "- Optimize operations.\n"
    "- Enhance marketing strategies.\n"
    "- Drive growth across ExpressMart locations.\n\n"
    "Implementing these recommendations will address underperformance and capitalize on strengths."
)

# Save the presentation
output_path = "C:/Users/USER/Documents/ExpressMart_Sales_Analysis_Report.pptx"

presentation.save(output_path)

output_path
